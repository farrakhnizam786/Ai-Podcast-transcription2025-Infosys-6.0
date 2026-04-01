from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # 1. Create Presentation Object
    prs = Presentation()

    # --- HELPER: Add Slide ---
    def add_slide(title_text, content_points):
        # Layout 1 is Title + Content
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title = slide.shapes.title
        title.text = title_text
        
        # Set Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = point
            p.font.size = Pt(18)
            p.space_after = Pt(10)

            # Check for sub-bullets (denoted by starting with "- ")
            if point.startswith("    -"):
                p.level = 1
                p.text = point.replace("    - ", "")

    # --- SLIDE 1: Title ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Podcast AI Analytics Dashboard"
    slide.placeholders[1].text = "Internship Project: Infosys\nAutomated Audio Analysis using Local AI"

    # --- SLIDE 2: Project Overview ---
    add_slide("1. Project Overview", [
        "Objective: Build an offline, privacy-focused tool for audio analysis.",
        "Problem Solved: Converts unstructured audio data into structured insights.",
        "Core Value Proposition:",
        "    - No Cloud Dependency (Privacy Safe)",
        "    - Zero API Costs (Runs locally)",
        "    - Multi-model Integration (Speech, Text, Sentiment)"
    ])

    # --- SLIDE 3: Technology Stack ---
    add_slide("2. Technology Stack", [
        "Language: Python 3.9+",
        "Frontend: Streamlit (Interactive Dashboard)",
        "AI Models:",
        "    - Transcription: OpenAI Whisper (Base model)",
        "    - Summarization: HuggingFace Transformers (DistilBART)",
        "    - Sentiment: NLTK VADER & TextBlob",
        "    - Topic Extraction: Scikit-Learn (TF-IDF & Cosine Similarity)",
        "Audio Processing: Librosa & SoundFile",
        "Visualization: Plotly Express"
    ])

    # --- SLIDE 4: Backend Architecture & Evolution ---
    add_slide("3. Backend Development", [
        "Initial State: Simple script processing files linearly.",
        "Optimization & Performance:",
        "    - Implemented Model Caching (Singleton Pattern) to reduce load times.",
        "File System Architecture:",
        "    - Structured storage: /audio, /transcripts, /summary, /sentiment_data.",
        "    - JSON-based persistence ensures data isn't lost on restart.",
        "Logic Refinements:",
        "    - Robust imports (fallback strategies).",
        "    - Audio format normalization (16kHz WAV conversion)."
    ])

    # --- SLIDE 5: Key Features (The Updates) ---
    add_slide("4. Feature Development", [
        "Smart Summarization: Generates abstractive summaries of episodes.",
        "Topic Segmentation:",
        "    - Mathematically detects conversation shifts.",
        "    - Includes Audio Player synced to specific segment timestamps.",
        "    - 'Download Info' button for saving segment details.",
        "Granular Sentiment Analysis:",
        "    - Calculates emotional tone (Positive/Negative/Neutral) per topic.",
        "    - Visualized via color-coded tags in the UI.",
        "Search & Highlight: Full-text search within transcripts."
    ])

    # --- SLIDE 6: Dashboard UI/UX ---
    add_slide("5. Dashboard UI Implementation", [
        "Framework: Streamlit",
        "Layout Structure:",
        "    - Sidebar: Uploads (File/URL) & Feature Toggles (Save resources).",
        "    - Tabs: Overview, Sentiment, Transcript.",
        "Visual Polish:",
        "    - Removed emojis/comments for professional look.",
        "    - Custom CSS for cards, metric boxes, and sentiment labels.",
        "    - Responsive plotting with Plotly."
    ])

    # --- SLIDE 7: Testing & Quality Assurance ---
    add_slide("6. Testing & Reliability", [
        "Unit Testing:",
        "    - Integrated 'unittest' for validating time-conversion logic.",
        "Error Handling:",
        "    - Try-Except blocks for missing files or failed downloads.",
        "    - User-friendly error messages in the UI.",
        "Robustness:",
        "    - Dynamic path handling (pathlib) for Windows compatibility.",
        "    - Checkbox toggles to control heavy AI processes."
    ])

    # --- SAVE FILE ---
    output_file = "Podcast_AI_Project_Report.pptx"
    prs.save(output_file)
    print(f"Success! Presentation saved as: {output_file}")

if __name__ == "__main__":
    create_presentation()